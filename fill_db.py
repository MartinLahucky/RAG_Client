import asyncio
import hashlib
import os
import time
from concurrent.futures import ThreadPoolExecutor
from typing import Dict, List, Any
import fitz
import magic
import olefile
import pandas as pd
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from lxml import etree
from pptx import Presentation
from watchdog.events import FileSystemEventHandler
from watchdog.observers import Observer
import logger
import tokenizer
import utils
from database import MongoDB


class NewFileHandler(FileSystemEventHandler):
    def __init__(self, process_function):
        self.process_function = process_function

    def on_created(self, event):
        if not event.is_directory:
            logger.log_info(f"Nový soubor detekován: {event.src_path}")
            self.process_function()


def monitor_directory(directory, process_function):
    event_handler = NewFileHandler(process_function)
    observer = Observer()
    observer.schedule(event_handler, directory, recursive=False)
    observer.start()
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()


def convert_metadata(metadata):
    converted = {}
    for key, value in metadata.items():
        if key in ['tokens', 'pos_tags', 'named_entities']:
            converted[key] = value
        elif isinstance(value, (list, tuple)):
            converted[key] = str(value)
        elif isinstance(value, (str, int, float, bool)):
            converted[key] = value
        else:
            converted[key] = str(value)
    return converted


def split_text(raw_documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=100,
        chunk_overlap=20,
        length_function=len,
        separators=[r"\n\n", r"\n", r"\.", r"!", r"\?", r"\d+\)", r"\d+\.", ]
    )

    if isinstance(raw_documents, list):
        if isinstance(raw_documents[0], dict):
            texts = [doc['page_content'] for doc in raw_documents]
        else:
            texts = raw_documents
    else:
        texts = [raw_documents]

    split_texts = []
    for text in texts:
        split_texts.extend(text_splitter.split_text(text))

    return [{"page_content": text} for text in split_texts]


def get_file_type(file_path):
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)

    # Fallback na detekci podle přípony
    if file_type == 'application/octet-stream':
        _, extension = os.path.splitext(file_path)
        extension = extension.lower()
        if extension == '.pdf':
            return 'application/pdf'
        elif extension in ['.doc', '.docx']:
            return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif extension in ['.xls', '.xlsx']:
            return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif extension in ['.ppt', '.pptx']:
            return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

    return file_type


def extract_text_from_pdf(file_path):
    try:
        document = fitz.open(file_path)
        text = ""
        for page in document:
            text += page.get_text()
        document.close()
        return [{"page_content": text}]  # Balení textu do slovníku
    except Exception as e:
        logger.log_warning(f"Error reading {file_path}: {e}")
        return []


def extract_text_from_ole_doc(file_path):
    try:
        # Kontrola, zda je soubor OLE formát
        if olefile.isOleFile(file_path):
            ole = olefile.OleFileIO(file_path)
            # Získáme obsah uložený v OLE souboru
            if ole.exists('WordDocument'):
                stream = ole.openstream('WordDocument')
                data = stream.read()
                # Pro jednoduchost zde vracíme pouze binární data; lze rozšířit
                text = data.decode('utf-8', 'ignore')  # Zkusíme dekódovat obsah
                if text.strip():
                    return [{"page_content": text}]
                else:
                    logger.log_warning(f"Extrahovaný text z OLE souboru {file_path} je prázdný.")
                    return []
            else:
                logger.log_warning(f"OLE soubor {file_path} neobsahuje WordDocument stream.")
                return []
        else:
            logger.log_warning(f"Soubor {file_path} není platný OLE formát.")
            return []
    except Exception as e:
        logger.log_warning(f"Chyba při extrakci textu z OLE souboru {file_path}: {str(e)}")
        return []


def extract_text_from_docx(file_path):
    try:
        # Zkusíme nejprve načíst jako Word dokument
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        if not text.strip():
            logger.log_warning(f"Extrahovaný text z dokumentu {file_path} je prázdný.")
            return []
        return [{"page_content": text}]
    except Exception as e:
        logger.log_warning(f"Chyba při extrakci textu z dokumentu {file_path}: {str(e)}")
        logger.log_info(f"Pokusíme se{file_path} načíst jako XML...")
        # Fallback - kontrola, zda není soubor XML
        try:
            with open(file_path, 'rb') as file:
                header = file.read(1024).decode('utf-8', 'ignore')
                if not header.strip().startswith('<?xml'):
                    logger.log_warning(f"Soubor {file_path} není platný XML formát.")
                    return []

                # Zkusíme načíst jako XML, pokud hlavička souboru naznačuje XML
                file.seek(0)  # Vrátíme čtecí hlavu na začátek souboru
                tree = etree.parse(file)
                text = ''.join(tree.xpath('//text()'))
                if text.strip():
                    return [{"page_content": text}]
                else:
                    logger.log_warning(f"Extrahovaný text z XML {file_path} je prázdný.")
        except Exception as xml_error:
            logger.log_warning(f"Chyba při extrakci textu z XML {file_path}: {str(xml_error)}")

        return []


import xlrd


def extract_text_from_xls(file_path):
    try:
        # Otevření souboru ve starším formátu Excelu
        workbook = xlrd.open_workbook(file_path)
        sheet = workbook.sheet_by_index(0)  # První list
        # Načteme všechny řádky a sloupce jako text
        text = '\n'.join([str(sheet.row_values(row)) for row in range(sheet.nrows)])
        if text.strip():
            return [{"page_content": text}]
        else:
            logger.log_warning(f"Extrahovaný text z Excelu {file_path} je prázdný.")
            return []
    except Exception as e:
        logger.log_warning(f"Chyba při extrakci textu z Excel {file_path}: {str(e)}")
        return []


def extract_text_from_xlsx(file_path):
    try:
        # Primárně načítáme pomocí pandas
        df = pd.read_excel(file_path, engine='openpyxl')
        text = df.to_string(index=False)
        return [{"page_content": text}]
    except Exception as e:
        logger.log_warning(f"Chyba při extrakci textu z Excel {file_path}: {str(e)}")
        logger.log_info(f"Pokusíme se{file_path} načíst jako starý formát pomocí xlrd...")
        # Fallback - pokusíme se načíst jako starý formát pomocí xlrd
        try:
            workbook = xlrd.open_workbook(file_path, on_demand=True)
            sheet = workbook.sheet_by_index(0)
            text = '\n'.join([str(sheet.row_values(row)) for row in range(sheet.nrows)])
            if text.strip():
                return [{"page_content": text}]
            else:
                logger.log_warning(f"Extrahovaný text z Excelu {file_path} je prázdný.")
        except Exception as xlrd_error:
            logger.log_warning(f"Chyba při extrakci textu ze starého Excel {file_path}: {str(xlrd_error)}")

        return []


def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        if not text:
            return [{"page_content": ""}]
        return [{"page_content": '\n'.join(text)}]
    except Exception as e:
        logger.log_warning(f"Chyba při extrakci textu z prezentace {file_path}: {str(e)}")
        return []


def extract_text_from_txt(file_path: str) -> List[Dict[str, str]]:
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return [{"page_content": content}]
    except Exception as e:
        logger.log_warning(f"Error reading text file {file_path}: {str(e)}")
        return []


def process_file(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    if not os.path.exists(file_path):
        logger.log_warning(f"File not found: {file_path}")
        return []

    try:
        # Kontrola duplicit
        file_hash = utils.calculate_file_hash(file_path)

        # Kontrola, zda dokument s tímto file_hash již existuje
        db = MongoDB()
        existing_doc = db.query_documents("data", {"metadata.file_hash": file_hash})
        if existing_doc:
            logger.log_info(f"Document {file_path} with hash {file_hash} already exists, skipping insertion.")
            return []

        # Zpracování PDF souborů
        if 'pdf' in file_type:
            return extract_text_from_pdf(file_path)

        # Zpracování souborů Word ve formátu DOCX
        elif 'wordprocessingml.document' in file_type or file_path.endswith('.docx'):
            return extract_text_from_docx(file_path)

        # Zpracování souborů Word ve starším formátu DOC (OLE)
        elif 'msword' in file_type or file_path.endswith('.doc'):
            return extract_text_from_ole_doc(file_path)

        # Zpracování souborů Excel ve formátu XLSX
        elif 'spreadsheetml.sheet' in file_type or file_path.endswith('.xlsx'):
            return extract_text_from_xlsx(file_path)

        # Zpracování souborů  ve starším formátu XLS
        elif 'ms-excel' in file_type or file_path.endswith('.xls'):
            return extract_text_from_xls(file_path)

        # Zpracování souborů PowerPoint
        elif 'ms-powerpoint' in file_type or 'presentationml.presentation' in file_type or file_path.endswith(
                ('.ppt', '.pptx')):
            return extract_text_from_pptx(file_path)

        # Zpracování prostého textu
        elif 'text/plain' in file_type:
            return extract_text_from_txt(file_path)

        # Nepodporovaný typ souboru
        else:
            logger.log_warning(f"Unsupported file type: {file_type} for file: {file_path}")
            return []

    except Exception as e:
        logger.log_warning(f"Error processing file {file_path}: {str(e)}")
        return []


def process_paragraph(paragraph, file_path):
    page_content = paragraph['page_content']
    tokens = tokenizer.tokenize_text(page_content)
    pos_tags = tokenizer.pos_tag(tokens)
    named_entities = tokenizer.named_entity_recognition(pos_tags)

    metadata = {
        "tokens": tokens,
        "pos_tags": pos_tags,
        "named_entities": named_entities,
        "source": file_path,
        "page": paragraph.get("page", 0)
    }

    converted_metadata = convert_metadata(metadata)

    # Vytvoření hash z obsahu dokumentu
    content_hash = hashlib.md5(page_content.encode('utf-8')).hexdigest()

    return {
        "_id": content_hash,
        "content": page_content,
        "metadata": converted_metadata
    }


async def process_document(file_path: str) -> List[Dict[str, Any]]:
    loop = asyncio.get_event_loop()
    file_type = await loop.run_in_executor(None, get_file_type, file_path)
    raw_documents = await loop.run_in_executor(None, process_file, file_path, file_type)
    if not raw_documents:
        logger.log_warning(f"Žádný obsah nebyl extrahován z {file_path}")
        return []

    paragraphs = await loop.run_in_executor(None, split_text, raw_documents)
    documents = await asyncio.gather(
        *[loop.run_in_executor(None, process_paragraph, paragraph, file_path) for paragraph in paragraphs])
    return documents


# Funkce pro načtení a zpracování dokumentů ve složce "data"
async def load_and_process_documents():
    logger.log_info("Zpracování dokumentů ve složce 'data'...")
    db = MongoDB()  # Připojení k MongoDB
    db.reload_localization()  # Načtení lokalizací

    documents = []
    file_paths = []
    for root, _, files in os.walk('data'):
        for file in files:
            file_paths.append(os.path.join(root, file))

    loop = asyncio.get_event_loop()
    with ThreadPoolExecutor(max_workers=4) as executor:
        tasks = [process_document(file_path) for file_path in file_paths]
        for future in asyncio.as_completed(tasks):
            try:
                result = await future
                documents.extend(result)
            except Exception as e:
                logger.log_warning(f"Chyba při zpracování souboru: {str(e)}")

    for document in documents:
        try:
            await loop.run_in_executor(None, db.insert_document, 'data', document)
        except Exception as e:
            logger.log_warning(f"Chyba při vkládání dokumentu: {str(e)}")

    await loop.run_in_executor(None, db.create_text_index, 'data', 'content')
    db.close_connection()
    logger.log_info("Dokumenty byly zpracovány a uloženy do databáze.")


if __name__ == "__main__":
    data_directory = 'data'
    if not os.path.exists(data_directory):
        os.makedirs(data_directory)
    monitor_directory(data_directory, load_and_process_documents)
